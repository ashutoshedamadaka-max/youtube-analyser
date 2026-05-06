import json
from pathlib import Path
from datetime import date
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DATA_DIR   = Path(".tmp/data")
CHARTS_DIR = Path(".tmp/charts")
THUMB_DIR  = Path(".tmp/thumbnails")
OUTPUT_DIR = Path("output")

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

NAVY   = RGBColor(0x1E, 0x2D, 0x4E)
TEAL   = RGBColor(0x4E, 0xCD, 0xC4)
CORAL  = RGBColor(0xFF, 0x6B, 0x6B)
YELLOW = RGBColor(0xFF, 0xE6, 0x6D)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LGREY  = RGBColor(0xF5, 0xF5, 0xF5)
MGREY  = RGBColor(0x99, 0x99, 0x99)


# ── helpers ──────────────────────────────────────────────────────────────────

def bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or WHITE


def rect(slide, l, t, w, h, fill_color, line_color=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if line_color:
        s.line.color.rgb = line_color
    else:
        s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=14, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or NAVY
    return tb


def bullets(slide, items, l, t, w, h, size=13):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(6)
        r = p.add_run()
        r.text = f"▸  {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = NAVY


def header(slide, title, color=None):
    rect(slide, 0, 0, SLIDE_W, Inches(1.05), color or NAVY)
    txt(slide, title, Inches(0.4), Inches(0.15), Inches(12.5), Inches(0.75),
        size=26, bold=True, color=WHITE)


def footer(slide, num):
    txt(slide, f"Confidential  •  {num}", Inches(0.3), Inches(7.15),
        Inches(12.6), Inches(0.25), size=9, color=MGREY, align=PP_ALIGN.RIGHT)


def img(slide, path, l, t, w, h):
    p = Path(path)
    if p.exists():
        slide.shapes.add_picture(str(p), l, t, w, h)


def divider(slide, t, color=None):
    rect(slide, Inches(0.5), t, Inches(12.33), Inches(0.04), color or TEAL)


def fmt(n):
    if n >= 1_000_000_000: return f"{n/1_000_000_000:.1f}B"
    if n >= 1_000_000:     return f"{n/1_000_000:.1f}M"
    if n >= 1_000:         return f"{n/1_000:.0f}K"
    return str(n)


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ── slides ────────────────────────────────────────────────────────────────────

def slide_title(prs, topic, today):
    s = blank(prs)
    bg(s, NAVY)
    rect(s, 0, Inches(3.15), SLIDE_W, Inches(0.08), TEAL)
    txt(s, "YouTube Space Analysis", Inches(1.5), Inches(1.4), Inches(10.33), Inches(0.7),
        size=18, color=TEAL, align=PP_ALIGN.CENTER)
    txt(s, topic.title(), Inches(0.8), Inches(2.05), Inches(11.73), Inches(1.4),
        size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, today, Inches(1), Inches(4.1), Inches(11.33), Inches(0.45),
        size=14, color=MGREY, align=PP_ALIGN.CENTER)
    txt(s, "Powered by YouTube Data API  +  Claude AI  +  python-pptx",
        Inches(1), Inches(6.6), Inches(11.33), Inches(0.4),
        size=10, color=MGREY, align=PP_ALIGN.CENTER)


def slide_exec_summary(prs, insights):
    s = blank(prs)
    bg(s)
    header(s, "Executive Summary")
    footer(s, 2)
    summ = insights.get("executive_summary", {})
    txt(s, summ.get("space_overview", ""), Inches(0.5), Inches(1.2), Inches(12.33), Inches(1.0), size=13)
    divider(s, Inches(2.35))
    bullets(s, summ.get("executive_summary_bullets", []),
            Inches(0.5), Inches(2.55), Inches(12.33), Inches(4.5), size=14)


def slide_top_videos(prs, videos):
    s = blank(prs)
    bg(s)
    header(s, "Top 10 Videos by Views")
    footer(s, 3)
    img(s, CHARTS_DIR / "top_videos.png", Inches(0.2), Inches(1.15), Inches(7.6), Inches(5.7))
    colors = [TEAL, CORAL, YELLOW]
    for i, v in enumerate(videos[:3]):
        t = Inches(1.4 + i * 1.85)
        rect(s, Inches(8.1), t, Inches(4.9), Inches(1.6), LGREY, colors[i])
        rect(s, Inches(8.1), t, Inches(0.18), Inches(1.6), colors[i])
        txt(s, f"#{i+1}", Inches(8.4), t + Inches(0.1), Inches(0.5), Inches(0.45),
            size=16, bold=True, color=colors[i])
        title_short = v["title"][:48] + ("…" if len(v["title"]) > 48 else "")
        txt(s, title_short, Inches(8.4), t + Inches(0.6), Inches(4.5), Inches(0.5), size=10)
        txt(s, f"{fmt(v['views'])} views  •  {v['channel_title'][:22]}",
            Inches(8.4), t + Inches(1.1), Inches(4.5), Inches(0.38),
            size=10, bold=True, color=CORAL)


def slide_top_channels(prs, channels):
    s = blank(prs)
    bg(s)
    header(s, "Top Channels in the Space")
    footer(s, 4)
    img(s, CHARTS_DIR / "top_channels.png", Inches(0.2), Inches(1.15), Inches(7.6), Inches(5.7))
    for i, c in enumerate(channels[:5]):
        t = Inches(1.45 + i * 1.1)
        txt(s, c["title"][:30], Inches(8.1), t, Inches(4.9), Inches(0.4), size=11, bold=True)
        txt(s, f"{fmt(c['subscribers'])} subs  •  {fmt(c['video_count'])} videos",
            Inches(8.1), t + Inches(0.4), Inches(4.9), Inches(0.35), size=10, color=MGREY)
        rect(s, Inches(8.1), t + Inches(0.82), Inches(4.6), Inches(0.02), LGREY)


def slide_engagement(prs):
    s = blank(prs)
    bg(s)
    header(s, "Engagement Deep Dive")
    footer(s, 5)
    img(s, CHARTS_DIR / "engagement.png", Inches(0.2), Inches(1.15), Inches(8.6), Inches(5.7))
    txt(s, "What this tells you:", Inches(9.15), Inches(1.4), Inches(3.8), Inches(0.4), size=13, bold=True)
    bullets(s, [
        "High views + low engagement = viral but not sticky",
        "High engagement + moderate views = loyal niche audience",
        "Target the upper-right: big reach AND resonance",
    ], Inches(9.15), Inches(1.9), Inches(3.8), Inches(2.5), size=11)
    img(s, CHARTS_DIR / "views_distribution.png", Inches(9.0), Inches(4.2), Inches(4.1), Inches(2.7))


def slide_content_themes(prs, insights):
    s = blank(prs)
    bg(s)
    header(s, "Content Themes")
    footer(s, 6)
    img(s, CHARTS_DIR / "content_themes.png", Inches(0.2), Inches(1.0), Inches(6.2), Inches(6.2))
    themes = insights.get("title_analysis", {}).get("dominant_content_themes", [])
    txt(s, "Theme Breakdown", Inches(6.7), Inches(1.4), Inches(6.3), Inches(0.4), size=13, bold=True)
    if themes:
        bullets(s, [f"{t['theme']}  ({t.get('percentage','?')}%)" for t in themes],
                Inches(6.7), Inches(1.9), Inches(6.3), Inches(3.8), size=12)


def slide_title_hooks(prs, insights):
    s = blank(prs)
    bg(s)
    header(s, "Title & Hook Patterns")
    footer(s, 7)
    ta = insights.get("title_analysis", {})
    txt(s, "Top Hook Patterns", Inches(0.5), Inches(1.2), Inches(5.8), Inches(0.4), size=13, bold=True)
    bullets(s, ta.get("top_hook_patterns", []),
            Inches(0.5), Inches(1.7), Inches(5.8), Inches(2.5), size=12)
    txt(s, "High-Performing Keywords", Inches(7.0), Inches(1.2), Inches(6.0), Inches(0.4), size=13, bold=True)
    bullets(s, ta.get("high_performing_keywords", [])[:8],
            Inches(7.0), Inches(1.7), Inches(6.0), Inches(2.5), size=12)
    divider(s, Inches(4.4))
    corr = ta.get("view_correlation_insights", "")
    if corr:
        txt(s, "Insight: " + corr, Inches(0.5), Inches(4.6), Inches(12.33), Inches(0.9), size=12)
    overused = ta.get("overused_patterns", [])
    if overused:
        txt(s, "⚠  Oversaturated — avoid:", Inches(0.5), Inches(5.65), Inches(5), Inches(0.4),
            size=12, bold=True, color=CORAL)
        txt(s, "  •  " + "   •  ".join(overused),
            Inches(0.5), Inches(6.1), Inches(12.33), Inches(0.45), size=11, color=MGREY)


def slide_thumbnail_trends(prs, insights, thumbnails):
    s = blank(prs)
    bg(s)
    header(s, "Thumbnail Visual Trends")
    footer(s, 8)
    shown = [t for t in thumbnails[:6] if Path(t["path"]).exists()]
    tw, th = Inches(2.55), Inches(1.55)
    sx, sy = Inches(0.35), Inches(1.25)
    for i, t in enumerate(shown):
        row, col = i // 3, i % 3
        try:
            slide_x = sx + col * (tw + Inches(0.15))
            slide_y = sy + row * (th + Inches(0.12))
            s.shapes.add_picture(t["path"], slide_x, slide_y, tw, th)
        except Exception:
            pass
    ta = insights.get("thumbnail_analysis", {})
    txt(s, "Visual Findings", Inches(8.3), Inches(1.35), Inches(4.7), Inches(0.4), size=13, bold=True)
    txt(s, ta.get("dominant_visual_style", ""), Inches(8.3), Inches(1.8), Inches(4.7), Inches(1.0), size=11)
    txt(s, "Color Trends", Inches(8.3), Inches(2.95), Inches(4.7), Inches(0.4), size=12, bold=True)
    bullets(s, ta.get("color_palette_trends", []), Inches(8.3), Inches(3.4), Inches(4.7), Inches(1.1), size=11)
    txt(s, "Design Recs", Inches(8.3), Inches(4.7), Inches(4.7), Inches(0.4), size=12, bold=True)
    bullets(s, ta.get("design_recommendations", []), Inches(8.3), Inches(5.15), Inches(4.7), Inches(1.8), size=11)


def slide_upload_cadence(prs, insights):
    s = blank(prs)
    bg(s)
    header(s, "Upload Cadence & Timing")
    footer(s, 9)
    img(s, CHARTS_DIR / "upload_cadence.png", Inches(0.2), Inches(1.15), Inches(8.6), Inches(5.7))
    posting = insights.get("executive_summary", {}).get("best_posting_insight", "")
    if posting:
        txt(s, "Posting Strategy Insight", Inches(9.15), Inches(1.4), Inches(3.9), Inches(0.4), size=13, bold=True)
        txt(s, posting, Inches(9.15), Inches(1.9), Inches(3.9), Inches(2.5), size=12)


def slide_content_gaps(prs, insights):
    s = blank(prs)
    bg(s)
    header(s, "Content Gaps & Opportunities")
    footer(s, 10)
    gaps = insights.get("executive_summary", {}).get("content_gaps", [])
    accents = [TEAL, CORAL, YELLOW]
    for i, gap in enumerate(gaps[:3]):
        t = Inches(1.55 + i * 1.85)
        rect(s, Inches(0.35), t, Inches(12.6), Inches(1.6), LGREY, accents[i])
        rect(s, Inches(0.35), t, Inches(0.18), Inches(1.6), accents[i])
        txt(s, f"Gap: {gap.get('gap', '')}", Inches(0.7), t + Inches(0.1),
            Inches(12.1), Inches(0.5), size=12, bold=True)
        txt(s, f"Opportunity: {gap.get('opportunity', '')}", Inches(0.7), t + Inches(0.65),
            Inches(12.1), Inches(0.75), size=11)


def slide_recommendations(prs, insights):
    s = blank(prs)
    bg(s)
    header(s, "Content Recommendations")
    footer(s, 11)
    recs = insights.get("executive_summary", {}).get("recommendations", [])[:5]
    accents = [TEAL, CORAL, YELLOW, NAVY, MGREY]
    cw = Inches(2.45)
    for i, rec in enumerate(recs):
        l = Inches(0.18) + i * (cw + Inches(0.1))
        t = Inches(1.25)
        rect(s, l, t, cw, Inches(5.9), LGREY, accents[i])
        rect(s, l, t, cw, Inches(0.2), accents[i])
        txt(s, f"#{i+1}", l + Inches(0.15), t + Inches(0.28), cw - Inches(0.2), Inches(0.4),
            size=18, bold=True, color=accents[i])
        txt(s, rec.get("title", ""), l + Inches(0.15), t + Inches(0.72),
            cw - Inches(0.2), Inches(0.8), size=11, bold=True)
        txt(s, rec.get("description", ""), l + Inches(0.15), t + Inches(1.6),
            cw - Inches(0.2), Inches(2.1), size=10)
        txt(s, "Why: " + rec.get("rationale", ""), l + Inches(0.15), t + Inches(3.85),
            cw - Inches(0.2), Inches(1.9), size=9, color=MGREY)


def slide_methodology(prs, videos, channels):
    s = blank(prs)
    bg(s)
    header(s, "Methodology & Data Sources")
    footer(s, 12)
    items = [
        f"Videos analyzed: {len(videos)}",
        f"Channels profiled: {len(channels)}",
        "Data source: YouTube Data API v3",
        "Transcript extraction: youtube-transcript-api (free)",
        "AI analysis: OpenAI GPT-4o (text) + GPT-4o-mini (vision)",
        "Chart generation: matplotlib",
        f"Report date: {date.today().strftime('%B %d, %Y')}",
    ]
    bullets(s, items, Inches(0.5), Inches(1.5), Inches(12.33), Inches(5.5), size=14)


# ── main ──────────────────────────────────────────────────────────────────────

def main():
    OUTPUT_DIR.mkdir(parents=True, exist_ok=True)

    with open(DATA_DIR / "videos.json",   encoding="utf-8") as f: vd       = json.load(f)
    with open(DATA_DIR / "channels.json", encoding="utf-8") as f: channels = json.load(f)["channels"]
    with open(DATA_DIR / "insights.json", encoding="utf-8") as f: insights = json.load(f)
    with open(DATA_DIR / "thumbnails.json", encoding="utf-8") as f: thumbs  = json.load(f)["thumbnails"]

    videos = vd["videos"]
    topic  = vd["topic"]
    today  = date.today().strftime("%B %d, %Y")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    steps = [
        ("Title",             lambda: slide_title(prs, topic, today)),
        ("Executive Summary", lambda: slide_exec_summary(prs, insights)),
        ("Top Videos",        lambda: slide_top_videos(prs, videos)),
        ("Top Channels",      lambda: slide_top_channels(prs, channels)),
        ("Engagement",        lambda: slide_engagement(prs)),
        ("Content Themes",    lambda: slide_content_themes(prs, insights)),
        ("Title Hooks",       lambda: slide_title_hooks(prs, insights)),
        ("Thumbnail Trends",  lambda: slide_thumbnail_trends(prs, insights, thumbs)),
        ("Upload Cadence",    lambda: slide_upload_cadence(prs, insights)),
        ("Content Gaps",      lambda: slide_content_gaps(prs, insights)),
        ("Recommendations",   lambda: slide_recommendations(prs, insights)),
        ("Methodology",       lambda: slide_methodology(prs, videos, channels)),
    ]

    for i, (name, fn) in enumerate(steps, 1):
        print(f"[{i:02d}/{len(steps)}] {name}")
        fn()

    slug = topic.replace(" ", "_")
    fname = f"{slug}_{date.today().strftime('%Y%m%d')}.pptx"
    out = OUTPUT_DIR / fname
    prs.save(str(out))
    print(f"\n✓ Saved → {out}")
    return str(out)


if __name__ == "__main__":
    main()
